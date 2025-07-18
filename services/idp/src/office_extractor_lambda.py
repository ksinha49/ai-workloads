# ------------------------------------------------------------------------------
# app.py
# ------------------------------------------------------------------------------
"""
Module: app.py
Description:
  Basic Lambda handler for 2-office-extractor.
Version: 1.0.0
"""

from __future__ import annotations
import io
import json
import os
from typing import Iterable
from pydantic import BaseModel, ValidationError
from models import S3Event, LambdaResponse

import boto3
from common_utils import get_config, configure_logger, iter_s3_records
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from ocr_module import convert_to_markdown

__author__ = "Koushik Sinha"
__version__ = "1.0.0"
__modified_by__ = "Koushik Sinha"

logger = configure_logger(__name__)

s3_client = boto3.client("s3")


class OfficeRecord(BaseModel):
    bucket: str
    key: str

    @classmethod
    def parse_record(cls, record: dict) -> "OfficeRecord":
        bucket = record.get("s3", {}).get("bucket", {}).get("name")
        key = record.get("s3", {}).get("object", {}).get("key")
        return cls(bucket=bucket, key=key)



def _extract_docx(body: bytes) -> list[str]:
    """Return Markdown pages extracted from a DOCX file."""

    doc = Document(io.BytesIO(body))
    lines = [p.text for p in doc.paragraphs if p.text]
    text = "\n".join(lines)
    return [convert_to_markdown(text, 1)]

def _extract_pptx(body: bytes) -> list[str]:
    """Return Markdown pages extracted from a PPTX file, one per slide."""

    pres = Presentation(io.BytesIO(body))
    pages: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        page_text = "\n".join(texts)
        pages.append(convert_to_markdown(page_text, i))
    return pages

def _extract_xlsx(body: bytes) -> list[str]:
    """Return Markdown pages extracted from an XLSX file, one per sheet."""

    wb = load_workbook(filename=io.BytesIO(body), data_only=True)
    pages: list[str] = []
    for i, sheet in enumerate(wb, start=1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if cell is None else str(cell) for cell in row]
            rows.append("| " + " | ".join(cells) + " |")
        text = "\n".join(rows)
        pages.append(convert_to_markdown(text, i))
    return pages

def _process_record(record: dict, document_id: str | None = None) -> None:
    """Extract text from an Office document referenced by ``record``.

    The object is read from S3, converted into Markdown pages depending on
    its file type and the result is stored as a JSON document under
    ``TEXT_DOC_PREFIX``.
    """

    try:
        rec = OfficeRecord.parse_record(record)
    except ValidationError as exc:
        logger.error("Invalid record: %s", exc)
        return
    bucket = rec.bucket
    key = rec.key
    bucket_name = get_config("BUCKET_NAME", bucket, key)
    office_prefix = get_config("OFFICE_PREFIX", bucket, key) or os.environ.get("OFFICE_PREFIX", "")
    text_doc_prefix = get_config("TEXT_DOC_PREFIX", bucket, key) or os.environ.get("TEXT_DOC_PREFIX")
    if office_prefix and not office_prefix.endswith("/"):
        office_prefix += "/"
    if text_doc_prefix and not text_doc_prefix.endswith("/"):
        text_doc_prefix += "/"
    if bucket != bucket_name or not key:
        logger.info("Skipping record with bucket=%s key=%s", bucket, key)
        return
    if not key.startswith(office_prefix):
        logger.info("Key %s outside prefix %s - skipping", key, office_prefix)
        return

    ext = os.path.splitext(key)[1].lower()
    if ext not in {".docx", ".pptx", ".xlsx"}:
        logger.info("Unsupported extension %s - skipping", ext)
        return

    obj = s3_client.get_object(Bucket=bucket_name, Key=key)
    body = obj["Body"].read()

    if ext == ".docx":
        pages = _extract_docx(body)
        typ = "docx"
    elif ext == ".pptx":
        pages = _extract_pptx(body)
        typ = "pptx"
    else:
        pages = _extract_xlsx(body)
        typ = "xlsx"

    doc_id = document_id or os.path.splitext(os.path.basename(key))[0]
    dest_key = f"{text_doc_prefix}{doc_id}.json"

    payload = {
        "documentId": doc_id,
        "type": typ,
        "pageCount": len(pages),
        "pages": pages,
    }

    s3_client.put_object(
        Bucket=bucket_name,
        Key=dest_key,
        Body=json.dumps(payload).encode("utf-8"),
        ContentType="application/json",
    )
    logger.info("Wrote %s", dest_key)

def lambda_handler(event: S3Event, context: dict) -> LambdaResponse:
    """Triggered by Office files arriving in ``OFFICE_PREFIX``.

    Parameters
    ----------
    event : :class:`models.S3Event`
        Standard S3 trigger event.

    1. Reads each S3 record and extracts text from DOCX, PPTX or XLSX files.
    2. Writes the resulting Markdown pages to ``TEXT_DOC_PREFIX``.

    Returns
    -------
    :class:`models.LambdaResponse`
        Response with a success message once processing completes.
    """

    logger.info("Received event for 2-office-extractor: %s", event)
    doc_id = getattr(event, "document_id", None)
    if isinstance(event, dict):
        doc_id = event.get("document_id", doc_id)
    for rec in iter_s3_records(event):
        try:
            _process_record(rec, doc_id)
        except Exception as exc:  # pragma: no cover - runtime safety
            logger.error("Error processing record %s: %s", rec, exc)
    return {
        "statusCode": 200,
        "body": json.dumps({"message": "2-office-extractor executed"})
    }
